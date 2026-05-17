"""Extract plain text from PDF/DOCX/PPTX for summarization."""
import io
import logging

logger = logging.getLogger(__name__)


def extract_text(data: bytes, ext: str) -> str:
    ext = (ext or "").lower()
    try:
        if ext == "pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            chunks = []
            for page in reader.pages[:80]:  # cap pages
                try:
                    chunks.append(page.extract_text() or "")
                except Exception:
                    continue
            return "\n\n".join(chunks)
        if ext in ("docx",):
            from docx import Document
            doc = Document(io.BytesIO(data))
            parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
            for tbl in doc.tables:
                for row in tbl.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells if cell.text))
            return "\n".join(parts)
        if ext in ("pptx",):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                lines = [f"--- Slide {i} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in para.runs)
                            if text.strip():
                                lines.append(text)
                slides.append("\n".join(lines))
            return "\n\n".join(slides)
        # Legacy .doc / .ppt — no clean extractor without external binaries
        return ""
    except Exception as e:
        logger.warning(f"Text extraction failed for ext={ext}: {e}")
        return ""
